"""Generate ClinVar benchmark slide deck (.pptx) for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x00, 0xC8, 0xAA)   # teal
ACCENT2   = RGBColor(0xFF, 0xC1, 0x07)   # amber
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GY  = RGBColor(0xCC, 0xCC, 0xDD)
DARK_TEXT = RGBColor(0x12, 0x12, 0x20)
ROW_EVEN  = RGBColor(0x22, 0x22, 0x40)
ROW_ODD   = RGBColor(0x2C, 0x2C, 0x50)
HEADER_BG = RGBColor(0x00, 0xA8, 0x8A)

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # dark background rectangle
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def accent_bar(slide, top=Inches(1.05), height=Inches(0.06)):
    bar = slide.shapes.add_shape(1, 0, top, W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    add_textbox(slide, title,
                left=Inches(0.5), top=Inches(0.1),
                width=Inches(12.3), height=Inches(0.9),
                font_size=28, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    left=Inches(0.5), top=Inches(1.15),
                    width=Inches(12.3), height=Inches(0.5),
                    font_size=14, color=LIGHT_GY)


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, highlight_col=None):
    """Draw a styled table manually using rectangles + textboxes."""
    n_cols = len(headers)
    n_rows = len(rows)

    if col_widths is None:
        cw = [width / n_cols] * n_cols
    else:
        cw = [Inches(w) for w in col_widths]

    row_h = height / (n_rows + 1)

    def cell(text, x, y, w, h, bg, fg=WHITE, bold=False, font_size=13, align=PP_ALIGN.CENTER):
        rect = slide.shapes.add_shape(1, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = DARK_BG
        rect.line.width = Pt(0.5)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    # header row
    x = left
    for i, h in enumerate(headers):
        cell(h, x, top, cw[i], row_h, HEADER_BG, WHITE, bold=True, font_size=13)
        x += cw[i]

    # data rows
    for r, row in enumerate(rows):
        x = left
        bg = ROW_EVEN if r % 2 == 0 else ROW_ODD
        for c, val in enumerate(row):
            fg = ACCENT2 if (highlight_col is not None and c == highlight_col) else LIGHT_GY
            bold = (c == highlight_col)
            cell(str(val), x, top + row_h * (r + 1), cw[c], row_h,
                 bg, fg, bold=bold, font_size=12)
            x += cw[c]


def bullet_box(slide, items, left, top, width, height, font_size=15, color=LIGHT_GY):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

# big teal title
add_textbox(slide, "vibe-vep ClinVar Benchmark",
            Inches(0.6), Inches(1.8), Inches(12), Inches(1.6),
            font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_textbox(slide, "Independent HGVSp Accuracy vs. snpEff & Ensembl VEP",
            Inches(0.6), Inches(3.5), Inches(12), Inches(0.7),
            font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, "232,008 pathogenic/likely-pathogenic variants  ·  GENCODE v49 / Ensembl 115",
            Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
            font_size=16, color=LIGHT_GY, align=PP_ALIGN.CENTER)

# bottom accent line
bar = slide.shapes.add_shape(1, 0, Inches(6.9), W, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT2; bar.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 – Methodology
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Methodology", "Ground truth independent of all annotation tools")

bullets = [
    "• Ground truth: ClinVar curated pathogenic/likely-pathogenic SNVs & indels",
    "  — Submitted by clinical labs, expert panels, and NCBI curators",
    "  — Completely independent of VEP, snpEff, or vibe-vep",
    "",
    "• 100% map to MANE Select transcripts (231,858 / 232,008)",
    "  — All tools use the same transcript → protein notation differences = real errors",
    "",
    "• Version-exact subset: 100% use the identical NM_ version (231,699 / 232,008)",
    "  — Transcript version drift cannot affect these results",
    "  — Failures are true algorithmic errors, not drift artifacts",
    "",
    "• Database versions: GENCODE v49 / Ensembl 115 (vibe-vep, VEP)  ·  snpEff GRCh38.115",
]
bullet_box(slide, bullets, Inches(0.6), Inches(1.7), Inches(12), Inches(5.3),
           font_size=16)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 – Overall Results
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Overall Results — HGVSp Match",
            "Primary transcript (best) and any transcript accuracy across 232,008 variants")

headers = ["Tool", "Best Match", "Any Transcript", "Not Annotated"]
rows = [
    ["vibe-vep",          "87.4%", "95.9%", "1,508"],
    ["snpEff GRCh38.115", "82.5%", "95.9%",   "318"],
    ["Ensembl VEP v115",  "79.8%", "96.2%",   "837"],
]
add_table(slide, headers, rows,
          left=Inches(1.5), top=Inches(1.8),
          width=Inches(10.3), height=Inches(2.2),
          col_widths=[3.5, 2.0, 2.5, 2.3],
          highlight_col=1)

add_textbox(slide,
            "vibe-vep leads on primary transcript selection  (+5% vs VEP, +5% vs snpEff)\n"
            "All three tools reach ~96% when any transcript is considered → gap is transcript choice, not algorithm.",
            Inches(0.6), Inches(4.4), Inches(12), Inches(1.5),
            font_size=15, color=ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 – By Variant Type
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "HGVSp Match by Variant Type", "Primary transcript (best) accuracy")

headers = ["Consequence Class", "n", "vibe-vep", "snpEff", "VEP"]
rows = [
    ["Missense",          "69,516", "91.3%", "86.0%", "80.8%"],
    ["Frameshift",        "83,151", "88.5%", "84.1%", "82.3%"],
    ["Stop-gained",       "75,508", "83.1%", "78.9%", "77.3%"],
    ["Inframe deletion",   "2,019", "86.8%", "82.1%", "78.3%"],
    ["Inframe insertion",    "748", "83.3%", "60.6%", "71.5%"],
    ["Synonymous",           "815", "89.8%",  "0.0%",  "0.0%"],
]
add_table(slide, headers, rows,
          left=Inches(0.6), top=Inches(1.7),
          width=Inches(12.1), height=Inches(3.6),
          col_widths=[2.8, 1.4, 2.0, 2.0, 2.0],
          highlight_col=2)

add_textbox(slide,
            "* vibe-vep uniquely annotates synonymous variants (e.g. p.Arg273=) — snpEff and VEP omit HGVSp for silent changes.",
            Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
            font_size=12, color=LIGHT_GY)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 – Failure Taxonomy
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Failure Taxonomy",
            "Three mutually exclusive failure types account for all mismatches")

headers = ["Failure Type", "vibe-vep", "snpEff", "VEP", "Root Cause"]
rows = [
    ["Any-not-best\n(transcript choice)",
     "8.4%", "13.4%", "16.5%",
     "Correct HGVSp in secondary transcript; ClinVar pre-dates MANE Select"],
    ["Complete fail\n(algorithmic)",
     "3.5%", "3.9%", "3.4%",
     "No transcript has the expected HGVSp; mostly ClinVar format artifacts"],
    ["Not annotated",
     "1,508", "318", "837",
     "Tool emits no protein change (splice/intron/UTR)"],
]
add_table(slide, headers, rows,
          left=Inches(0.3), top=Inches(1.7),
          width=Inches(12.7), height=Inches(2.8),
          col_widths=[2.4, 1.2, 1.2, 1.2, 6.7],
          highlight_col=1)

add_textbox(slide,
            "Key insight: vibe-vep's MANE Select preference minimises transcript-choice errors (8.4% vs 16.5% for VEP).\n"
            "Complete-failure rates are nearly identical across tools — dominated by ClinVar format artifacts, not algorithmic differences.",
            Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.4),
            font_size=15, color=ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 – Known Gaps & Fixes
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Known Gaps & Actionable Fixes",
            "Remaining vibe-vep-specific failures are fixable engineering issues")

headers = ["Issue", "Consequence", "Impact", "Fix"]
rows = [
    ["Splice-adjacent frameshifts\nemit _splice suffix HGVSp",
     "frameshift", "~0.6% miss\n(85% vibe-specific)",
     "Improve splice/frameshift\npriority at exon boundaries"],
    ["Stop-creating inframe deletions\nmisclassified as stop_gained",
     "inframe_del", "~6.1% miss\n(96% vibe-specific)",
     "Preserve inframe_del class\nwhen reading frame is intact"],
    ["ClinVar p.Met1? vs p.Met1Arg\nformat inconsistency",
     "missense", "2.5% shared\n(all tools equal)",
     "ClinVar artifact — not fixable;\nshared baseline failure"],
    ["ClinVar p.Xaa_YbbinsTer\nrange notation",
     "stop_gained", "~70% of stop_gained\ncomplete-fail (shared)",
     "ClinVar artifact — not fixable;\nshared baseline failure"],
]
add_table(slide, headers, rows,
          left=Inches(0.3), top=Inches(1.7),
          width=Inches(12.7), height=Inches(4.1),
          col_widths=[3.5, 2.0, 2.4, 4.8])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 – Performance
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Performance", "232,008 variants annotated end-to-end")

headers = ["Tool", "Throughput", "Wall Time", "Speed-up vs VEP"]
rows = [
    ["vibe-vep",          "23,722 v/s", "9.8 s",   "136×"],
    ["snpEff GRCh38.115",    "513 v/s", "452 s",     "3×"],
    ["Ensembl VEP v115",     "174 v/s", "1,333 s",   "—"],
]
add_table(slide, headers, rows,
          left=Inches(1.2), top=Inches(1.8),
          width=Inches(10.9), height=Inches(2.0),
          col_widths=[3.2, 2.5, 2.2, 3.0],
          highlight_col=1)

add_textbox(slide,
            "vibe-vep is 46× faster than snpEff and 136× faster than VEP\n"
            "Cache warm-up: 2.0 s from DuckDB binary cache  ·  snpEff/VEP times from *.elapsed sidecar files",
            Inches(0.6), Inches(4.2), Inches(12), Inches(1.0),
            font_size=15, color=ACCENT2)

# Big numbers
for label, val, xpos in [
    ("23,722 v/s", "vibe-vep", Inches(1.5)),
    ("46×", "vs snpEff", Inches(5.5)),
    ("136×", "vs VEP", Inches(9.0)),
]:
    add_textbox(slide, label,
                xpos, Inches(5.2), Inches(3), Inches(0.8),
                font_size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, val,
                xpos, Inches(5.9), Inches(3), Inches(0.4),
                font_size=14, color=LIGHT_GY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 – Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_title(slide, "Summary")

# 4 stat boxes
boxes = [
    ("87.4%",  "Best HGVSp accuracy\n(+5% vs VEP)",          Inches(0.4)),
    ("97.7%",  "Consequence class\naccuracy (highest)",        Inches(3.6)),
    ("136×",   "Faster than Ensembl VEP\n(DuckDB cache)",      Inches(6.8)),
    ("2 fixes","Close most remaining\nvibe-specific gaps",     Inches(10.0)),
]
for val, label, x in boxes:
    rect = slide.shapes.add_shape(1, x, Inches(1.7), Inches(2.9), Inches(1.9))
    rect.fill.solid(); rect.fill.fore_color.rgb = ROW_ODD
    rect.line.color.rgb = ACCENT; rect.line.width = Pt(1.5)
    add_textbox(slide, val,   x, Inches(1.75), Inches(2.9), Inches(0.9),
                font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x, Inches(2.65), Inches(2.9), Inches(0.8),
                font_size=13, color=LIGHT_GY, align=PP_ALIGN.CENTER)

bullets = [
    "• Best primary HGVSp accuracy across all three tools on 232k independent ClinVar variants",
    "• Highest consequence class accuracy at 97.7%",
    "• Unique: correctly annotates synonymous variants (snpEff and VEP emit nothing)",
    "• Orders-of-magnitude faster with DuckDB binary caching",
    "• Remaining gaps are fixable: splice/frameshift boundary logic, inframe-del stop classification",
    "• ClinVar format artifacts (p.Met1?, range notation) account for most shared failures — not algorithmic",
]
bullet_box(slide, bullets, Inches(0.6), Inches(3.85), Inches(12), Inches(3.0),
           font_size=16)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "testdata/clinvar/benchmark_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
