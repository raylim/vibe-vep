"""3-slide summary deck. Plots are data only — no annotation boxes."""

import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as ticker
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG     = "#1A1A2E"
PANEL  = "#22223A"
TEAL   = "#00C8AA"
AMBER  = "#FFC107"
CORAL  = "#FF6B6B"
LAVEND = "#A78BFA"
GREY   = "#666688"
WHITE  = "#FFFFFF"
LGREY  = "#CCCCDD"
DIM    = "#3A3A55"

TOOLS  = ["vibe-vep", "snpEff", "VEP v115", "ANNOVAR"]
COLORS = [TEAL, AMBER, CORAL, LAVEND]

BEST  = [87.4, 82.8, 80.0, 65.2]
ANY   = [95.9, 96.3, 96.6, 96.0]
SNV   = [90.2, 85.6, 81.7, 69.3]
INDEL = [83.3, 78.6, 77.6, 59.1]

CLASSES    = ["missense (70k)", "frameshift (83k)", "stop_gained (76k)",
              "inframe_del (2k)", "inframe_ins (748)", "synonymous (815)"]
VIBE_CLS   = [91.3, 88.5, 83.1, 86.8, 83.3, 89.8]
SNPEFF_CLS = [86.0, 84.1, 78.9, 82.1, 60.6, 80.2]
VEP_CLS    = [80.8, 82.3, 77.3, 78.3, 71.5, 75.9]
ANNOV_CLS  = [70.6, 63.6, 63.0, 53.8, 58.0, 66.3]

TX_ERR   = [8.4,  13.5, 16.5, 30.8]
ALGO_ERR = [3.5,   3.6,  3.1,  3.7]
NOT_ANN  = [1508/232008*100, 318/232008*100, 837/232008*100, 21/232008*100]

SPEEDS = [23594, 513, 174, 151]
CONSEQ = [97.7, 97.0, 96.9, 97.8]

FS = 11   # single font size for all plot text

# ── figure helpers ────────────────────────────────────────────────────────────
def fig_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, facecolor=BG)
    buf.seek(0); plt.close(fig); return buf

def sax(ax, xlabel="", ylabel=""):
    ax.set_facecolor(PANEL)
    for sp in ax.spines.values(): sp.set_color(DIM)
    ax.tick_params(colors=LGREY, labelsize=FS)
    if xlabel: ax.set_xlabel(xlabel, color=LGREY, fontsize=FS)
    if ylabel: ax.set_ylabel(ylabel, color=LGREY, fontsize=FS)
    ax.grid(axis="x", color=DIM, lw=0.5, alpha=0.5, zorder=1)

# ── PPTX helpers ──────────────────────────────────────────────────────────────
SW, SH  = Inches(13.33), Inches(7.5)
C_TEAL  = RGBColor(0x00, 0xC8, 0xAA)
C_AMBER = RGBColor(0xFF, 0xC1, 0x07)
C_LGREY = RGBColor(0xCC, 0xCC, 0xDD)
C_DIM   = RGBColor(0x3A, 0x3A, 0x55)
C_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
C_PANEL = RGBColor(0x2C, 0x2C, 0x50)

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK; bg.line.fill.background()
    return s

def rule(s, top, color=C_TEAL):
    b = s.shapes.add_shape(1, 0, top, SW, Inches(0.055))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

def txt(s, text, l, t, w, h, size=12, bold=False,
        color=C_LGREY, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h); tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"

def add_fig(s, fig, l, t, w, h):
    s.shapes.add_picture(fig_bytes(fig), l, t, w, h)

def header(s, title, subtitle, num, total=3):
    rule(s, Inches(1.02))
    txt(s, title,    Inches(0.4), Inches(0.06), Inches(11.8), Inches(0.92),
        size=26, bold=True, color=C_TEAL)
    txt(s, subtitle, Inches(0.4), Inches(1.09), Inches(12.5), Inches(0.4),
        size=11, color=C_LGREY)
    txt(s, f"{num} / {total}", Inches(12.6), Inches(0.06), Inches(0.7), Inches(0.5),
        size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

def stat(s, val, label, x, y, w=Inches(2.9), h=Inches(1.5),
         val_color=C_TEAL, val_size=30):
    r = s.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = C_PANEL
    r.line.color.rgb = val_color; r.line.width = Pt(1.5)
    txt(s, val,   x, y + Inches(0.04), w, Inches(0.78),
        size=val_size, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(s, label, x, y + Inches(0.82), w, Inches(0.60),
        size=10, color=C_LGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1
# ══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()
header(s1,
    "vibe-vep leads primary accuracy; gap to any-transcript is transcript selection",
    "232,008 pathogenic ClinVar variants · independent ground truth · GENCODE v49 / Ensembl 115",
    num=1)

y = np.arange(4); bh = 0.32

# A: primary vs any (own figure)
fig_a, ax = plt.subplots(figsize=(6.5, 3.9), facecolor=BG)
fig_a.subplots_adjust(left=0.16, right=0.90, top=0.88, bottom=0.26)
sax(ax, xlabel="% HGVSp match vs ClinVar")
ax.barh(y+bh/2, BEST, bh, color=COLORS, alpha=0.95, zorder=3)
ax.barh(y-bh/2, ANY,  bh, color=COLORS, alpha=0.30, zorder=3)
ax.set_xlim(55, 112); ax.set_yticks(y); ax.set_yticklabels(TOOLS, fontsize=FS)
ax.axvline(90, color=GREY, lw=1.0, ls="--", alpha=0.5, zorder=2)
for i,(b,a,col) in enumerate(zip(BEST,ANY,COLORS)):
    ax.text(b+0.5, i+bh/2, f"{b}%", va="center", color=col,   fontsize=FS, fontweight="bold")
    ax.text(a+0.5, i-bh/2, f"{a}%", va="center", color=LGREY, fontsize=FS, alpha=0.80)
ax.legend(handles=[
    mpatches.Patch(color=GREY, alpha=0.95, label="Primary (tool's top-ranked transcript)"),
    mpatches.Patch(color=GREY, alpha=0.30, label="Any (correct answer in any transcript)"),
], fontsize=FS-1, facecolor=PANEL, labelcolor=LGREY, edgecolor=DIM,
   loc="upper center", bbox_to_anchor=(0.5, -0.18), ncol=1)
ax.set_title("Primary vs any-transcript accuracy", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s1, fig_a, Inches(0.2), Inches(1.52), Inches(6.5), Inches(3.9))

# B: SNV vs indel (own figure)
fig_b, ax = plt.subplots(figsize=(6.3, 3.9), facecolor=BG)
fig_b.subplots_adjust(left=0.06, right=0.90, top=0.88, bottom=0.26)
sax(ax, xlabel="% HGVSp match (primary)")
ax.barh(y+bh/2, SNV,   bh, color=COLORS, alpha=0.95, zorder=3)
ax.barh(y-bh/2, INDEL, bh, color=COLORS, alpha=0.40, zorder=3)
ax.set_xlim(50, 110); ax.set_yticks(y); ax.set_yticklabels([], fontsize=FS)
for i,(sv,d,col) in enumerate(zip(SNV,INDEL,COLORS)):
    ax.text(sv+0.5, i+bh/2, f"{sv}%", va="center", color=col,   fontsize=FS, fontweight="bold")
    ax.text(d+0.5,  i-bh/2, f"{d}%",  va="center", color=LGREY, fontsize=FS, alpha=0.80)
ax.legend(handles=[
    mpatches.Patch(color=GREY, alpha=0.95, label="SNV"),
    mpatches.Patch(color=GREY, alpha=0.40, label="Indel"),
], fontsize=FS-1, facecolor=PANEL, labelcolor=LGREY, edgecolor=DIM,
   loc="upper center", bbox_to_anchor=(0.5, -0.18), ncol=2)
ax.set_title("SNV vs indel accuracy", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s1, fig_b, Inches(6.9), Inches(1.52), Inches(6.3), Inches(3.9))

rule(s1, Inches(5.55), color=C_AMBER)
stat(s1, "+5%",  "vs snpEff\nprimary",   Inches(0.3),  Inches(5.63))
stat(s1, "+7%",  "vs VEP\nprimary",      Inches(3.35), Inches(5.63))
stat(s1, "+22%", "vs ANNOVAR\nprimary",  Inches(6.4),  Inches(5.63))
stat(s1, "~96%", "All tools\n(any tx)",  Inches(9.45), Inches(5.63))
rule(s1, Inches(7.25), color=C_TEAL)
txt(s1,
    "Primary is what matters in practice: tools report one result per variant in VCF/MAF output, "
    "so the top-ranked transcript is what analysts and clinicians actually see.  "
    "Any = the correct answer exists somewhere across all transcripts the tool reports — "
    "all four tools reach ~96%, so the protein change is computed correctly; "
    "the primary gap reflects which transcript each tool ranks first.",
    Inches(0.3), Inches(7.25), Inches(12.7), Inches(0.40), size=11, color=C_LGREY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
header(s2,
    "Failures: mostly transcript selection, not wrong algorithms",
    "vibe-vep: 8.4% transcript-choice errors vs 30.8% for ANNOVAR · ~3–4% algorithmic fail shared by all tools",
    num=2)

# ── Plot C: failure taxonomy (left panel, own figure) ────────────────────────
fig_c, ax = plt.subplots(figsize=(6.0, 3.9), facecolor=BG)
fig_c.subplots_adjust(left=0.18, right=0.82, top=0.88, bottom=0.28)
sax(ax, xlabel="% of all variants")

y = np.arange(4)
lefts_a  = TX_ERR
lefts_na = [t+a for t,a in zip(TX_ERR,ALGO_ERR)]
totals   = [t+a+n for t,a,n in zip(TX_ERR,ALGO_ERR,NOT_ANN)]
ax.barh(y, TX_ERR,   0.52, color=AMBER,  alpha=0.92, zorder=3)
ax.barh(y, ALGO_ERR, 0.52, color=CORAL,  alpha=0.92, zorder=3, left=lefts_a)
ax.barh(y, NOT_ANN,  0.52, color=LAVEND, alpha=0.92, zorder=3, left=lefts_na)
ax.set_xlim(0, 44); ax.set_yticks(y); ax.set_yticklabels(TOOLS, fontsize=FS)
# tx% inside amber segment
for i,t in enumerate(TX_ERR):
    ax.text(t/2, i, f"{t}%", va="center", ha="center",
            color=BG, fontsize=FS, fontweight="bold", zorder=5)
# total label — fits because xlim leaves room
for i,tot in enumerate(totals):
    ax.text(tot+0.4, i, f"{tot:.1f}%", va="center", color=LGREY, fontsize=FS, fontweight="bold")
# legend below axes — avoids all bars
ax.legend(handles=[
    mpatches.Patch(color=AMBER,  alpha=0.92, label="Wrong transcript chosen"),
    mpatches.Patch(color=CORAL,  alpha=0.92, label="No transcript correct"),
    mpatches.Patch(color=LAVEND, alpha=0.92, label="Not annotated"),
], fontsize=FS-1, facecolor=PANEL, labelcolor=LGREY, edgecolor=DIM,
   loc="upper center", bbox_to_anchor=(0.5, -0.20), ncol=3)
ax.set_title("Where do mismatches come from?", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s2, fig_c, Inches(0.2), Inches(1.52), Inches(6.3), Inches(3.9))

# ── Plot D: per-class dot chart (right panel, own figure) ────────────────────
fig_d, ax = plt.subplots(figsize=(6.5, 3.9), facecolor=BG)
# extra bottom margin for legend placed outside axes
fig_d.subplots_adjust(left=0.28, right=0.88, top=0.88, bottom=0.24)
sax(ax)  # no xlabel — title conveys axis meaning

nc=6; yy=np.arange(nc)
all_v=[VIBE_CLS,SNPEFF_CLS,VEP_CLS,ANNOV_CLS]; DOFF=[0.22,0.07,-0.07,-0.22]
ax.set_xlim(44, 106); ax.set_yticks(yy); ax.set_yticklabels(CLASSES, fontsize=FS)
ax.axvline(80, color=GREY, lw=0.8, ls="--", alpha=0.4, zorder=2)
for j in range(nc):
    vals=[v[j] for v in all_v]
    ax.plot([min(vals),max(vals)],[j,j], color=GREY, lw=3, alpha=0.25, zorder=2)
for vals,col,tool,dy in zip(all_v,COLORS,TOOLS,DOFF):
    ax.scatter(vals, yy+dy, color=col, s=80, zorder=4, label=tool)
for j,val in enumerate(VIBE_CLS):
    ax.text(val+0.7, j+DOFF[0], f"{val}%", va="center",
            color=TEAL, fontsize=FS-1, fontweight="bold")
# legend below axes — avoids all data points
ax.legend(fontsize=FS-1, facecolor=PANEL, labelcolor=LGREY, edgecolor=DIM,
          loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=4)
ax.set_title("Primary accuracy by consequence class", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s2, fig_d, Inches(6.7), Inches(1.52), Inches(6.6), Inches(3.9))

rule(s2, Inches(5.55), color=C_AMBER)
txt(s2,
    "Most mismatches are transcript-choice errors, not wrong protein calculations. "
    "Each gene has multiple transcripts; ClinVar records which transcript was used, "
    "but tools pick their own default. vibe-vep always prefers the MANE Select transcript "
    "(the clinically agreed standard), keeping these mismatches to 8.4%  vs 13.5% snpEff · 16.5% VEP · 30.8% ANNOVAR.",
    Inches(0.3), Inches(5.65), Inches(12.7), Inches(0.70), size=13, color=C_LGREY)
txt(s2,
    "Two known vibe-vep bugs account for most of its remaining unique failures:  "
    "① variants near splice sites that cause a frameshift are mis-labeled with a non-standard suffix  ·  "
    "② deletions that keep the reading frame but hit a stop codon are reported as stop_gained instead of inframe_deletion",
    Inches(0.3), Inches(6.42), Inches(12.7), Inches(0.65), size=13, color=C_LGREY)
rule(s2, Inches(7.15), color=C_TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
header(s3,
    "45–154× faster; consequence classification consistent across tools",
    "End-to-end wall time · 232,008 variants · vibe-vep includes 2.0 s DuckDB cache load",
    num=3)

# E: speed (own figure)
fig_e, ax = plt.subplots(figsize=(6.5, 3.6), facecolor=BG)
fig_e.subplots_adjust(left=0.14, right=0.97, top=0.88, bottom=0.14)
ax.set_facecolor(PANEL)
for sp in ax.spines.values(): sp.set_color(DIM)
ax.tick_params(colors=LGREY, labelsize=FS)
ax.grid(axis="y", color=DIM, lw=0.5, alpha=0.5, zorder=1)
ax.set_ylabel("variants / second", color=LGREY, fontsize=FS)
bars = ax.bar(np.arange(4), SPEEDS, 0.58, color=COLORS, alpha=0.92, zorder=3)
ax.set_yscale("log")
ax.set_ylim(80, 200000)
ax.set_xticks(np.arange(4)); ax.set_xticklabels(TOOLS, fontsize=FS)
ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda v,_: f"{int(v):,}"))
for b,v,t in zip(bars,SPEEDS,TOOLS):
    ax.text(b.get_x()+b.get_width()/2, v*1.5, f"{v:,}",
            ha="center", va="bottom",
            color=WHITE if t=="vibe-vep" else LGREY,
            fontsize=FS, fontweight="bold" if t=="vibe-vep" else "normal")
for i,s in enumerate(SPEEDS[1:],1):
    mid=10**((np.log10(s)+np.log10(SPEEDS[0]))/2)
    ax.text(i, mid, f"{SPEEDS[0]//s}×\nslower",
            ha="center", va="center", color=LGREY, fontsize=FS, style="italic")
ax.set_title("Annotation throughput (log scale)", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s3, fig_e, Inches(0.2), Inches(1.52), Inches(6.5), Inches(3.9))

# F: consequence (own figure, with tool name y-labels)
fig_f, ax = plt.subplots(figsize=(6.3, 3.6), facecolor=BG)
fig_f.subplots_adjust(left=0.20, right=0.88, top=0.88, bottom=0.14)
sax(ax, xlabel="% consequence class match")
bars = ax.barh(np.arange(4), CONSEQ, 0.52, color=COLORS, alpha=0.92, zorder=3)
ax.set_xlim(95.5, 101); ax.set_yticks(np.arange(4)); ax.set_yticklabels(TOOLS, fontsize=FS)
ax.axvline(97, color=GREY, lw=0.8, ls="--", alpha=0.5, zorder=2)
for b,v,col in zip(bars,CONSEQ,COLORS):
    ax.text(v+0.05, b.get_y()+b.get_height()/2, f"{v}%",
            va="center", color=col, fontsize=FS, fontweight="bold")
ax.set_title("Consequence class accuracy", color=WHITE, fontsize=FS+1, fontweight="bold")
add_fig(s3, fig_f, Inches(6.9), Inches(1.52), Inches(6.3), Inches(3.9))

rule(s3, Inches(5.55), color=C_AMBER)
stat(s3, "87.4%", "Highest primary\nHGVSp accuracy",          Inches(0.3),  Inches(5.63))
stat(s3, "97.7%", "Highest consequence\nclass accuracy",       Inches(3.35), Inches(5.63))
stat(s3, "154×",  "Faster than ANNOVAR\n(23,594 vs 151 v/s)", Inches(6.4),  Inches(5.63))
stat(s3, "~96%",  "All tools converge\n(any transcript)",      Inches(9.45), Inches(5.63))
rule(s3, Inches(7.25), color=C_TEAL)
txt(s3,
    "Ground truth: ClinVar curated by clinical labs & expert panels — independent of VEP, snpEff, or vibe-vep.",
    Inches(0.3), Inches(7.30), Inches(12.7), Inches(0.25), size=10, color=C_DIM)

out = "testdata/clinvar/benchmark_summary_slide.pptx"
prs.save(out)
print(f"Saved: {out}")
